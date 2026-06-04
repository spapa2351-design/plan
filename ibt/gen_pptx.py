# -*- coding: utf-8 -*-
"""스타벅스 바리스타 평가 제안 — PT 슬라이드(pptx) 생성.
스토리라인: 문제 → 해결 → 데모(실화면 캡처) → 효과.
색=터치클래스(Toss) 톤 + 스타벅스 그린 액센트. 폰트=Malgun Gothic."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

GREEN      = RGBColor(0x00,0x70,0x4A)
GREEN_DK   = RGBColor(0x00,0x5C,0x3C)
GREEN_SOFT = RGBColor(0xE6,0xF2,0xED)
INK        = RGBColor(0x19,0x1F,0x28)
GREY       = RGBColor(0x6B,0x76,0x84)
GREY_LT    = RGBColor(0x8B,0x95,0xA1)
LINE       = RGBColor(0xD1,0xD6,0xDB)
LIGHT      = RGBColor(0xF3,0xF4,0xF6)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
WARN       = RGBColor(0xC2,0x70,0x0A)
WARN_SOFT  = RGBColor(0xFD,0xF1,0xE2)
BLUE       = RGBColor(0x1B,0x64,0xDA)
BLUE_SOFT  = RGBColor(0xE8,0xF3,0xFF)

FONT = "Malgun Gothic"
HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "shots")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide():
    return prs.slides.add_slide(BLANK)

def _set_font(run, size, bold, color, name=FONT):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', name)

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
    tb = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    if runs and isinstance(runs[0], tuple): runs=[runs]
    for i,line in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp
        for (txt,size,bold,color) in line:
            r = p.add_run(); r.text = txt; _set_font(r,size,bold,color)
    return tb

def box(s, l, t, w, h, fill=WHITE, border=LINE, bw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06):
    sp = s.shapes.add_shape(shape, Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = border; sp.line.width = Pt(bw)
    sp.shadow.inherit = False
    try:
        if shape==MSO_SHAPE.ROUNDED_RECTANGLE:
            sp.adjustments[0]=radius
    except Exception: pass
    sp.text_frame.clear()
    return sp

def card(s, l, t, w, h, title, lines=None, fill=WHITE, border=LINE, tcolor=INK, tsize=14, lsize=11, lcolor=GREY, sp=1.05):
    box(s,l,t,w,h,fill=fill,border=border)
    runs=[[(title,tsize,True,tcolor)]] if title else []
    if lines:
        for ln in lines:
            runs.append([(ln,lsize,False,lcolor)])
    if runs:
        text(s,l+0.18,t+0.13,w-0.32,h-0.24,runs,anchor=MSO_ANCHOR.TOP,sp=sp)

def arrow(s, l, t, w, h=0.30, fill=GREEN, shape=MSO_SHAPE.RIGHT_ARROW):
    sp = s.shapes.add_shape(shape, Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background()
    sp.shadow.inherit=False
    return sp

def chip(s, l, t, txt, fill=GREEN_SOFT, fg=GREEN_DK, w=None, border=None):
    w = w or (0.26+0.108*len(txt))
    box(s,l,t,w,0.32,fill=fill,border=border,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    text(s,l,t+0.03,w,0.27,[[(txt,10.5,True,fg)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    return w

def pic(s, name, l, t, w):
    p = os.path.join(SHOTS, name)
    sp = s.shapes.add_picture(p, Inches(l), Inches(t), width=Inches(w))
    sp.line.color.rgb = LINE; sp.line.width = Pt(1.0)
    sp.shadow.inherit = False
    return sp

def header(s, kicker, title, kc=GREEN):
    box(s,0,0,SW,1.18,fill=WHITE,border=None)
    box(s,0,1.18,SW,0.02,fill=LINE,border=None)
    text(s,0.6,0.22,12,0.3,[[(kicker,11,True,kc)]])
    text(s,0.6,0.5,12.2,0.6,[[(title,22,True,INK)]])

def footer(s, n):
    text(s,0.6,SH-0.42,9,0.3,[[("스타벅스 바리스타 평가 · 2026-06-04",9,False,GREY_LT)]])
    text(s,SW-1.2,SH-0.42,0.6,0.3,[[(str(n),9,False,GREY_LT)]],align=PP_ALIGN.RIGHT)

def bg(s, color=WHITE):
    box(s,0,0,SW,SH,fill=color,border=None,shape=MSO_SHAPE.RECTANGLE)

# 데모 슬라이드 우측 포인트(작은 카드 4개)
def points(s, items, l=8.4, t=1.55, w=4.4, ch=1.18, gap=0.07):
    y=t
    for it in items:
        if len(it)==3:
            title, body, kind = it
        else:
            title, body = it; kind="n"
        fill, border, tc = WHITE, LINE, INK
        if kind=="g": fill, border, tc = GREEN_SOFT, GREEN, GREEN_DK
        if kind=="w": fill, border, tc = WARN_SOFT, WARN, WARN
        card(s,l,y,w,ch,title,[body],fill=fill,border=border,tcolor=tc,tsize=13.5,lsize=11.5,lcolor=(INK if kind=="w" else GREY))
        y += ch+gap

# ═══════════════════════ 1. 표지 ═══════════════════════
s=slide(); bg(s,INK)
box(s,0.9,2.55,0.62,0.62,fill=GREEN,border=None,shape=MSO_SHAPE.OVAL)
text(s,0.88,3.5,11,0.4,[[("STARBUCKS · 바리스타 교육 평가",13,True,RGBColor(0x8F,0xD3,0xB8))]])
text(s,0.85,3.95,11.8,1.7,[
    [("매장에서 직접 해보고, 30분 안에 제출하면,",30,True,WHITE)],
    [("AI가 1차 채점하고 사람이 확정한다",30,True,WHITE)],
], sp=1.12)
text(s,0.85,5.6,11.8,0.5,[[("터치클래스 기반 온라인 수행·서술형 평가",14,False,RGBColor(0xB0,0xB8,0xC1))]])
text(s,0.85,6.6,11.8,0.4,[[("2026-06-04",11,False,GREY_LT)]])

# ═══════════════════════ 2. 문제 ═══════════════════════
s=slide(); bg(s); header(s,"문제 · PROBLEM","지금, 바리스타 실무 역량은 평가하기 어렵다", kc=WARN)
probs=[("정답이 하나가 아니다",
        ["테이스팅·센서리 표현은 객관식으로",
         "잡히지 않는 수행·서술 역량입니다.",
         "향·산미·바디를 말로 풀어내는 답은",
         "여러 갈래라, 채점 기준이 흐릿하고",
         "보는 사람마다 다르게 읽힙니다."]),
       ("한자리에 모이기 어렵다",
        ["전국 매장·교대 근무 구조에서",
         "바리스타를 한곳에 모아 평가하는 건",
         "운영 부담이 큽니다.",
         "각자 자기 매장에서 평가받을 수 있는",
         "원격 방식이 필요합니다."]),
       ("채점 부담이 크다",
        ["백지에서 서술형 답안을 읽고",
         "점수를 매기는 일은 채점자 부담이",
         "크고, 채점자마다 점수가 갈립니다.",
         "공정성과 일관성을 지키기가",
         "쉽지 않습니다."])]
x=0.6; w=3.97
for i,(t1,body) in enumerate(probs):
    box(s,x,1.55,w,4.75,fill=WHITE,border=LINE)
    text(s,x+0.28,1.85,0.6,0.6,[[(str(i+1),26,True,WARN)]])
    text(s,x+0.95,1.99,w-1.1,0.7,[[(t1,16,True,INK)]])
    text(s,x+0.28,2.85,w-0.5,3.3,[[(ln,12.5,False,GREY)] for ln in body],sp=1.35)
    x+=4.18
footer(s,2)

# ═══════════════════════ 3. 해결 한 줄 ═══════════════════════
s=slide(); bg(s); header(s,"해결 · SOLUTION","검증된 평가 엔진 위에 얹어, 빠르고 안전하게 도입한다")
box(s,0.6,1.55,12.13,1.85,fill=GREEN_SOFT,border=GREEN)
text(s,0.95,1.72,11.4,0.4,[[("한 줄 콘셉트",12,True,GREEN_DK)]])
text(s,0.95,2.12,11.6,1.2,[
    [("스타벅스가 자기 매장에서 직접 수행한 실습을 ",19,True,INK),("30분 제한",19,True,GREEN_DK),(" 안에 제출하면,",19,True,INK)],
    [("AI가 1차 채점",19,True,GREEN_DK),("하고 ",19,True,INK),("사람이 확정",19,True,GREEN_DK),("하는 온라인 평가.",19,True,INK)],
],sp=1.18)
text(s,0.6,3.7,12,0.4,[[("이 방식이 맞는 이유",13,True,GREEN_DK)]])
whys=[("수행·서술 역량","정답이 여러 갈래라 객관식이 아니라 루브릭 채점이 맞습니다"),
      ("내부 자료 활용","새 평가틀을 만들지 않고 스타벅스 내부 자료를 지문으로 그대로 녹입니다"),
      ("도입이 가볍다","검증된 엔진 기반이라 도입 리스크·비용 부담이 작습니다")]
x=0.6; w=3.97
for (t1,b) in whys:
    card(s,x,4.15,w,2.3,t1,None,fill=WHITE,border=LINE,tcolor=GREEN_DK,tsize=14)
    text(s,x+0.2,4.8,w-0.36,1.5,[[(b,12.5,False,GREY)]],sp=1.25)
    x+=4.18
footer(s,3)

# ═══════════════════════ 4. 해결 메커니즘 3축 ═══════════════════════
s=slide(); bg(s); header(s,"해결 · SOLUTION","세 가지 장치로, 앞의 세 가지 문제를 푼다")
mech=[("내부 자료 내재화",
       ["스타벅스 내부 자료를 평가 지문이자",
        "AI 채점 사전으로 함께 쓴다."],
       "→ 문제 ① 완화", "채점 어휘가 공통 사전으로 고정돼 기준이 또렷해진다", GREEN),
      ("개인 원격 30분",
       ["시간창 안에 각자 매장에서 30분 수행·제출.",
        "끊겨도 자동저장으로 복구. 행위는 시스템 밖."],
       "→ 문제 ② 해결", "모이지 않고 자기 매장에서 응시한다", BLUE),
      ("AI 1차 → 사람 확정",
       ["루브릭을 배점화 → AI가 충족판정·잠정점수·",
        "코멘트 초안. 사람이 손대야 확정."],
       "→ 문제 ③ 완화", "백지 채점이 아니라 초안 검수 + 점수차 누적 개선", GREEN)]
x=0.6; w=3.97
for i,(t1,body,tag,tagsub,c) in enumerate(mech):
    box(s,x,1.55,w,4.55,fill=WHITE,border=LINE)
    chip(s,x+0.25,1.78,"장치 "+str(i+1),fill=(BLUE_SOFT if c==BLUE else GREEN_SOFT),fg=(BLUE if c==BLUE else GREEN_DK))
    text(s,x+0.25,2.3,w-0.45,0.5,[[(t1,16,True,INK)]])
    text(s,x+0.25,3.0,w-0.45,1.3,[[(ln,12,False,GREY)] for ln in body],sp=1.2)
    box(s,x+0.25,4.35,w-0.5,1.5,fill=(BLUE_SOFT if c==BLUE else GREEN_SOFT),border=None)
    text(s,x+0.45,4.5,w-0.85,0.45,[[(tag,13.5,True,(BLUE if c==BLUE else GREEN_DK))]])
    text(s,x+0.45,4.95,w-0.85,0.85,[[(tagsub,11.5,False,INK)]],sp=1.15)
    x+=4.18
footer(s,4)

# ═══════════════════════ 5. 데모 전체 흐름 ═══════════════════════
s=slide(); bg(s); header(s,"데모 · DEMO","한 번 따라가 보겠습니다 — 만들고, 응시하고, 채점한다")
# 5노드 흐름
nodes=[("① 문제 빌더","관리자 · 만들기",GREEN_SOFT,GREEN_DK),
       ("② 30분 응시","학습자 · 직접 수행",GREEN_SOFT,GREEN_DK),
       ("제출","텍스트 + 사진",WHITE,INK),
       ("③ 채점·확정","관리자 · AI→사람",GREEN_SOFT,GREEN_DK),
       ("결과·이력","합격 = 자격 기록",WHITE,INK)]
y=2.55; w=2.1; h=1.5; step=2.5075
for i,(t1,b,fc,tc) in enumerate(nodes):
    x=0.6+step*i
    box(s,x,y,w,h,fill=fc,border=(GREEN if fc==GREEN_SOFT else LINE))
    text(s,x+0.16,y+0.26,w-0.3,0.5,[[(t1,14,True,tc)]])
    text(s,x+0.16,y+0.85,w-0.3,0.5,[[(b,11,False,GREY)]],sp=1.05)
    if i<4:
        arrow(s,x+w+0.05,y+h/2-0.15,0.31,fill=(INK if i==2 else GREEN))
# 레인 라벨
text(s,0.6,2.1,6,0.3,[[("학습자 ↔ 관리자가 같은 한 흐름 위에서 움직입니다",12,True,GREEN_DK)]])
box(s,0.6,4.45,12.13,1.0,fill=GREEN_SOFT,border=None)
text(s,0.85,4.78,12,0.5,[
    [("다음 3장에서 ",13,True,GREEN_DK),("실제 화면 캡처",13,True,GREEN_DK),("로 ① 빌더 · ② 응시 · ③ 채점·확정을 차례로 보여드립니다.",13,False,INK)],
])
footer(s,5)

# ═══════════════════════ 6. 데모 ① 빌더 ═══════════════════════
s=slide(); bg(s); header(s,"데모 ① · 관리자","문제 빌더 — 서술형 문항에 자료·가이드·루브릭을 얹는다")
pic(s,"builder.png",0.55,1.5,7.55)
points(s,[
    ("서술형 수행과제 문항","자료 첨부 · 작성가이드(단계) · 제한시간을 갖춘 문항"),
    ("자료 = 내부 자료","첨부한 내부 자료가 문제 지문·채점 사전과 연결된다","g"),
    ("루브릭에 배점","체크포인트마다 점수 → AI가 읽고 1차 채점"),
    ("합격·채점자 설정","70점 컷오프 · 점장 1인 / 2인 교차 토글"),
])
footer(s,6)

# ═══════════════════════ 7. 데모 ② 응시 ═══════════════════════
s=slide(); bg(s); header(s,"데모 ② · 학습자","30분 응시 — 자기 매장에서 직접 수행하고 서술 제출")
pic(s,"exam.png",0.55,1.5,7.55)
points(s,[
    ("30분 타이머","시간창 안에 시작 → 5분 전 경고 · 0초 자동제출"),
    ("좌 문제 / 우 답안","문제·작성가이드·내부 자료 + 답안 한 칸 + 사진"),
    ("끊겨도 복구","자동저장이 모바일 이탈을 보완 (PC 권장)"),
    ("'행위'는 시스템 밖","실제 시음·추출은 응시자가 직접. 시스템은 열람·제출·채점만","w"),
])
footer(s,7)

# ═══════════════════════ 8. 데모 ③ 채점·확정 ═══════════════════════
s=slide(); bg(s); header(s,"데모 ③ · 관리자","AI가 먼저 채점하고, 사람이 확정한다")
pic(s,"grading.png",0.55,1.5,7.55)
points(s,[
    ("3분할 한눈에","좌 루브릭·모범답안 / 중 답안 / 우 AI 1차"),
    ("AI 1차 (잠정 82)","향미 어휘 감지 · 루브릭 배점 합산 · 코멘트 초안","g"),
    ("사람이 확정 (85)","검수·조정해 확정. 자동 합·불 없음 — 사람이 손대야 확정"),
    ("자격 이력에 기록","확정 시 응시자 수강 이력에 자격 결과로 남는다"),
])
footer(s,8)

# ═══════════════════════ 9. 효과 셀링포인트 ═══════════════════════
s=slide(); bg(s); header(s,"효과 · EFFECT","그래서 무엇이 좋아지나")
sp3=[("내부 자료 활용",
      ["스타벅스 내부 자료를 평가 지문이자",
       "채점 어휘로 그대로 쓴다.",
       "새 평가틀을 만들 필요가 없다."]),
     ("채점 부담 경감",
      ["백지에서 채점하지 않고",
       "AI 1차 초안을 손보기만 한다.",
       "채점자 간 편차도 줄어든다."]),
     ("합격 = 자격 이력",
      ["통과 결과가 수강 이력에",
       "자격으로 남는다.",
       "이력 관리가 자연스럽게 이어진다."])]
x=0.6; w=3.97
for (t1,body) in sp3:
    box(s,x,1.55,w,4.6,fill=GREEN_SOFT,border=GREEN)
    text(s,x+0.28,1.95,w-0.5,0.6,[[(t1,17,True,GREEN_DK)]])
    box(s,x+0.28,2.7,1.1,0.02,fill=GREEN,border=None)
    text(s,x+0.28,2.95,w-0.5,2.9,[[(ln,12.5,False,INK)] for ln in body],sp=1.4)
    x+=4.18
footer(s,9)

# ═══════════════════════ 10. 범위 & 단계 ═══════════════════════
s=slide(); bg(s); header(s,"효과 · EFFECT","1차 제공 범위와 다음 단계")
# 좌: 단계별 제출 형식
text(s,0.6,1.55,7,0.4,[[("단계별로 여는 제출 형식",13,True,GREEN_DK)]])
box(s,0.6,2.05,5.97,1.95,fill=GREEN_SOFT,border=GREEN)
chip(s,0.85,2.27,"1단계 · 지금 제공",fill=GREEN,fg=WHITE)
text(s,0.85,2.95,5.5,0.5,[[("텍스트 서술 + 사진",16,True,GREEN_DK)]])
text(s,0.85,3.45,5.5,0.5,[[("결과물·단계 인증 사진까지 받습니다",12.5,False,INK)]])
box(s,0.6,4.15,5.97,1.95,fill=LIGHT,border=LINE)
chip(s,0.85,4.37,"2단계 · 이어서 확장",fill=GREY,fg=WHITE)
text(s,0.85,5.05,5.5,0.5,[[("음성 · 영상 · 수치 데이터",16,True,GREY)]])
text(s,0.85,5.55,5.5,0.5,[[("효과가 큰 순서로 단계적으로 엽니다",12.5,False,GREY)]])
# 우: 이번에 다루지 않는 것
box(s,6.78,2.05,5.95,4.05,fill=WHITE,border=LINE)
text(s,7.03,2.28,5.4,0.4,[[("이번 평가가 다루지 않는 것",14.5,True,INK)]])
notin=["오프라인 집합 교육·출석 관리","실시간 화면 감독(프록터링)","시스템 자동 합격·불합격 판정"]
yy=3.05
for it in notin:
    text(s,7.03,yy,5.4,0.5,[[("·  ",13,True,GREY_LT),(it,13,False,INK)]],sp=1.0)
    yy+=0.65
text(s,7.03,5.35,5.4,0.7,[[("평가 범위를 또렷이 해, 꼭 필요한 것에 집중합니다.",11.5,False,GREY)]],sp=1.25)
footer(s,10)

# ═══════════════════════ 11. 마무리 ═══════════════════════
s=slide(); bg(s); header(s,"마무리 · NEXT","미팅에서 함께 확정할 것", kc=WARN)
text(s,0.6,1.5,12,0.4,[[("운영 방식은 스타벅스와 함께 정합니다 — 아래를 논의 주제로 제안합니다.",13,True,INK)]])
qs=[("1  합격 기준","70점 컷오프가 맞나? 자격별로 기준이 다른가?"),
    ("2  재응시","기본 1회 + 관리자 부여로 충분한가?"),
    ("3  채점자","점장 1인인가, 2인 교차가 필요한 자격이 있나?"),
    ("4  제출 형식","사진까지가 1차로 적절한가? 영상·음성의 우선순위는?"),
    ("5  평가 범위","테이스팅 외에 어떤 실습까지 평가하나?")]
y=2.15
for (t1,b) in qs:
    box(s,0.6,y,12.13,0.82,fill=WHITE,border=LINE)
    text(s,0.9,y+0.04,1.7,0.74,[[(t1,15,True,WARN)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,2.6,y+0.04,9.7,0.74,[[(b,13,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=0.94
footer(s,11)

# ───── 저장 (열려 있으면 대체 파일명) ─────
primary = os.path.join(HERE, "스타벅스-바리스타-평가-제안.pptx")
try:
    prs.save(primary)
    print("saved:", primary, "slides:", len(prs.slides._sldIdLst))
except PermissionError:
    alt = os.path.join(HERE, "스타벅스-바리스타-평가-제안_v2.pptx")
    prs.save(alt)
    print("PRIMARY LOCKED - saved alt:", alt, "slides:", len(prs.slides._sldIdLst))
